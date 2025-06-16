import os
import re
import time
from typing import List, Dict, Any, Optional
from PIL import Image
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from .utils import logger, ensure_directory_exists

# Initialize Gemini client
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "AIzaSyDyHokKwcdbJ-iY0CKvQoglq7Bqyr7HIFM")
client = genai.Client(api_key=GEMINI_API_KEY)

def summarize_image_with_gemini(image_path: str, caption: str) -> str:
    """Summarize an image using Gemini."""
    with Image.open(image_path) as image:
        prompt = (
            f"The following image is from a scientific paper, and the caption if existed is:\n\n"
            f"{caption}\n\n"
            "Based on both the visual content and the caption if existed, write a detailed but concise description of the image. "
            "Keep the description within 100 words maximum. Use a neutral academic tone.\n\n"
            "Then, generate a very short image caption (one single sentence, 5–15 words) for displaying below the image in a slide. "
            "The short caption must be included even if you are unsure. Do not skip it.\n\n"
            "Return the result in the following format exactly:\n"
            "**Image Description**: <your detailed description here>\n"
            "**Image Caption**: <your short caption here>"
        )

        response = client.models.generate_content(
            model="gemini-2.0-flash",
            contents=[image, prompt]
        )
        return response.text.strip()

def summarize_table_with_gemini(caption: str, table_html: str) -> str:
    """Summarize a table using Gemini."""
    prompt = (
        f"You are an assistant tasked with analyzing and summarizing tables from scientific papers. "
        f"The following is an HTML table and its caption (if existed).\n\n"
        f"Table Caption: {caption}\n"
        f"Table Data (HTML):\n{table_html}\n\n"
        "Analyze the table and provide a concise summary of the main results, explicitly including key numerical values. "
        "Highlight which methods performed better on each dataset (TotalText and CTW-1500) based on recall (R), precision (P), and H-mean (H). "
        "When identifying the best-performing methods, include the actual metric values (e.g., 'CRAFT achieved the highest recall of 85.3% on CTW-1500'). "
        "Use a clear and objective academic tone. Do not restate the caption.\n\n"
        "After the description, generate a a **very short, concise table caption (one single sentence, 5–15 words) that clearly summarizes the key insight of the table.\n\n"
        "Return the result in the following format:\n"
        "**Table Description**: <your detailed table analysis here>\n"
        "**Table Caption**: <your concise caption here>"
    )
    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=[prompt]
    )
    return response.text.strip()

def prompt_ver1(text: str, equations: List[Dict[str, str]]) -> str:
    """Generate prompt for version 1 summarization."""
    eq_list = "\n".join([
        f"- {eq.get('equation_markdown', '')} (filename: {eq.get('src', '')}"
        for eq in equations
    ]) if equations else "[None]"

    return f"""
        You are a scientific assistant helping to summarize and organize technical research documents.

        Your task:

        1. Summarize the content into **no more than 5 short bullet points**:
        • Each bullet must start with the • symbol (U+2022).
        • Use plain English and write **short, concise, clear sentences**.
        • Focus only on the **main contributions, methods, or findings**.
        • Avoid redundancy. Each point should convey a unique idea.
        • Only write as many bullets as needed — do not force 5 bullets if the content is limited.

        2. Analyze the equations listed below:
        • If a bullet point refers to or is supported by an equation, add the filename in this format: (Equation: filename.png) at the end of that bullet.
        • Do NOT include math notation or equation content in the bullets.
        • Only annotate relevant equations.

        3. Make sure the summary reflects all key insights from the text and equations.
        • Do NOT refer to images or tables explicitly.
        • Use plain English. Avoid any LaTeX, symbols, or code.

        **Output Format (strict):**
        • Bullet 1  
        • Bullet 2  
        (up to 5 bullets)

        Do NOT include any titles, explanations, or extra text.

        ---

        **Raw Text:**
        {text.strip()}

        ---

        **Equations:**
        {eq_list}

        Write concise summary bullets now, annotated with equations where appropriate.
        """

def prompt_ver2(text: str, equations: List[Dict[str, str]], images: List[Dict[str, str]]) -> str:
    """Generate prompt for version 2 summarization."""
    eq_list = "\n".join([
        f"- {eq.get('equation_markdown', '')} (filename: {eq.get('src', '')}"
        for eq in equations
    ]) if equations else "[None]"

    max_section = len(images) + 1

    figure_list = []
    table_list = []
    for img in images:
        alt = img.get("alt", "")
        name = img.get("src", "")
        caption = img.get("caption", "").strip()
        item_str = f"(Caption: {caption}) - (Filename: {name})"
        if alt.lower() == "table":
            table_list.append(item_str)
        else:
            figure_list.append(item_str)

    img_list = "\n".join(figure_list) if figure_list else "[None]"
    table_list_str = "\n".join(table_list) if table_list else "[None]"

    return f"""
        You are a scientific assistant helping to summarize and organize technical research documents.

        Your task is as follows:

        **Input:**
        - A raw text block from a scientific paper (see below).
        - A list of related equations in LaTeX/Markdown format (with filenames).
        - A list of related figures and tables (each with a caption and filename).

        ---

        **Your Goals:**

        1. **Summarize** the input text into concise, clear bullet points.
        2. **Organize** the bullets into logical sections:
        - Each section should have a header: `### Section N: Title of section`
        - Each section can have **up to 5 bullets**.
        - Each section can reference **at most 2 visuals** (images or tables).
        - The **total number of sections must not exceed {max_section}**.

        3. **Annotate** each bullet with any clearly related visual or equation:
        - Use format: `(Image: figure1.png, figure2.jpg) or (Table: table1.png) or (Equation: 1.png)`
        - Use at most 2 references per bullet.
        - If no relevant visual/equation applies, leave it untagged.
        

        **Coverage Requirement:**

        - Make sure that **each image, table, and equation is referenced by at least one bullet** if relevant.
        - Distribute references across sections so that **no image/table/equation is left out unnecessarily**.
        - Avoid forcing unrelated visuals into bullets — only include them where they naturally fit.


        **Formatting Guidelines:**
        - Use natural language to explain all ideas — do NOT include any mathematical notation, LaTeX, or special symbols (like $...$)
        - Each bullet must be a standalone factual summary.
        - Do not reuse or rephrase raw text — rewrite clearly.
        - Do not invent filenames or captions.
        - Do not add more than {max_section} sections.


        **Raw Text:**
        {text.strip()}

        ---

        **Equations:**
        {eq_list}

        ---

        **Figures:**
        {img_list}

        ---

        **Tables:**
        {table_list_str}

        ---

        Now write the summary bullets, organized by section, and annotate related visuals/equations accordingly.
        """

def summarize_block(block: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Summarize a single block using either prompt_ver1 or prompt_ver2."""
    text = block['text'].strip()
    if not text:
        logger.info(f"Block ('{block['heading']}') is empty, skipping.")
        return []

    logger.info(f"Summarizing block — {block['heading'][:30]}...")
    block_summaries = []

    num_images = len(block.get('images', []))
    num_equations = len(block.get('equations', []))
    
    try:
        if num_images > 2 or num_equations > 3:
            prompt = prompt_ver2(
                text=text,
                equations=block.get("equations", []),
                images=block.get("images", [])
            )
            response = client.models.generate_content(
                model="gemini-2.0-flash",
                contents=prompt
            )
            summarized = response.text.strip()
            
            sections = parse_summary_into_sections(summarized, block)
            block_summaries.extend(sections)
            
        else:
            prompt = prompt_ver1(
                text=text,
                equations=block.get("equations", [])
            )
            
            response = client.models.generate_content(
                model="gemini-2.0-flash", 
                contents=prompt
            )
            time.sleep(2)
            summarized = response.text.strip()
            
            summary_block = {
                "heading": block["heading"],
                "text": summarized,
                "equations": block.get("equations", []),
                "images": block.get("images", [])
            }
            block_summaries.append(summary_block)

        logger.info("Done summarizing block")
        return block_summaries

    except Exception as e:
        logger.error(f"Error summarizing block: {e}")
        return [{
            "heading": block["heading"],
            "text": "[Summary failed]",
            "equations": block.get("equations", []),
            "images": block.get("images", [])
        }]

def parse_summary_into_sections(summarized_text: str, block: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Parse summarized text into sections."""
    sections = []
    current_section = None
    current_bullets = []
    image_filenames = set()

    image_lookup = {img["src"]: img for img in block.get("images", [])}

    lines = summarized_text.strip().splitlines()
    heading = block["heading"]

    for line in lines:
        line = line.strip()
        if line.startswith("### Section"):
            if current_section and current_bullets:
                sections.append({
                    "heading": heading,
                    "sub_heading": current_section,
                    "text": "\n".join(current_bullets).strip(),
                    "images": [
                        image_lookup[fname]
                        for fname in sorted(image_filenames)
                        if fname in image_lookup
                    ]
                })

            current_section = line.replace("###", "").strip()
            current_bullets = []
            image_filenames = set()

        elif line.startswith("*"):
            content = line[1:].strip()

            related_imgs = re.findall(r"\(Image: ([^)]+)\)", content)
            related_tables = re.findall(r"\(Table: ([^)]+)\)", content)
            image_filenames.update(related_imgs + related_tables)

            clean_content = re.sub(r"\(Image: [^)]+\)", "", content)
            clean_content = re.sub(r"\(Table: [^)]+\)", "", clean_content).strip()

            current_bullets.append(clean_content)

    if current_section and current_bullets:
        sections.append({
            "heading": heading,
            "sub_heading": current_section,
            "text": "\n".join(current_bullets).strip(),
            "images": [
                image_lookup[fname]
                for fname in sorted(image_filenames)
                if fname in image_lookup
            ]
        })

    return sections

def create_presentation(blocks: List[Dict[str, Any]], pdf_filename: str, output_file: str = "output.pptx") -> None:
    """Create PowerPoint presentation from blocks."""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    def add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        slide.placeholders[1].text_frame.word_wrap = True

    def add_vertical_layout_slide(title: str, text: str) -> None:
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes

        title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.clear()
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True

        left = Inches(0.5)
        top = Inches(1)
        width = Inches(8.5)
        add_bullets_with_equations(shapes, left, top, width, text)

    def add_image_slide(title: str, text: str, images: List[Dict[str, str]]) -> None:
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes

        title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.clear()
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True

        left = Inches(0.5)
        top = Inches(1)
        width = Inches(4.3)
        add_bullets_with_equations(shapes, left, top, width, text, max_image_width=Inches(3.5))

        y = 1
        for img in images:
            path = img['path']
            caption = img.get('caption', '')
            if not os.path.exists(path):
                logger.warning(f"Image not found: {path}")
                continue

            pic = slide.shapes.add_picture(path, Inches(5.3), Inches(y), width=Inches(4))
            pic_height = pic.height.inches

            if caption:
                caption_top = y + pic_height + 0.1
                cap_box = slide.shapes.add_textbox(Inches(5.3), Inches(caption_top), Inches(4), Inches(0.6))
                cap_tf = cap_box.text_frame
                cap_tf.clear()
                cap_tf.paragraphs[0].text = caption
                cap_tf.paragraphs[0].font.size = Pt(12)

            y = y + pic_height + 0.5

    for i, block in enumerate(blocks):
        title = block.get('heading', '')
        text = block.get('text', '')
        images = []

        if i == 0:
            add_title_slide(title, text)
        elif i == 1:
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            tf = title_box.text_frame
            tf.clear()
            tf.paragraphs[0].text = title
            tf.paragraphs[0].font.size = Pt(36)
            tf.paragraphs[0].font.bold = True

            content_box = shapes.add_textbox(Inches(0.5), Inches(1), Inches(8.5), Inches(5))
            set_text_with_font(content_box.text_frame, text, 24)
        else:
            if block.get('images'):
                for img in block['images']:
                    filename = os.path.basename(img['src'])
                    full_path = os.path.join("/kaggle/working", pdf_filename, "auto/images", filename)
                    caption = img.get('caption', '')
                    images.append({'path': full_path, 'caption': caption})

            if images:
                add_image_slide(title, text, images)
            else:
                add_vertical_layout_slide(title, text)

    ensure_directory_exists(output_file)
    prs.save(output_file)
    logger.info(f"PowerPoint file created: {output_file}")

def set_text_with_font(text_frame, text: str, font_size: int = 16, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    """Set text with specified font and alignment."""
    text_frame.word_wrap = True
    text_frame.clear()
    for line in text.strip().splitlines():
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.alignment = align

def add_bullets_with_equations(shapes, left: float, top: float, width: float, text: str, max_image_width: float = Inches(3.5)) -> float:
    """Add bullets with equations to slide."""
    eq_pattern = re.compile(r'\(Equation:\s*([^\)]+)\)')
    current_top = top
    bullet_height = Inches(1)
    spacing = Inches(0.2)

    bullets = [line.strip() for line in text.split('\n') if line.strip()]
    for bullet in bullets:
        eq_match = eq_pattern.search(bullet)
        eq_path = None
        if eq_match:
            eq_path = eq_match.group(1).strip()
            bullet = eq_pattern.sub('', bullet).strip()

        text_box = shapes.add_textbox(left, current_top, width, bullet_height)
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.clear()

        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)

        current_top += bullet_height + spacing

        if eq_path and os.path.exists(eq_path):
            pic = shapes.add_picture(eq_path, left + (width - max_image_width) / 2, current_top, width=max_image_width)
            pic_height = pic.height.inches
            current_top += Inches(pic_height) + spacing
        elif eq_path:
            logger.warning(f"Equation image not found: {eq_path}")

    return current_top 